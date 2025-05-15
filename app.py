import ast
import comtypes.client
import configparser
import cv2
import copy
import dspy
import io
import logging
import numpy as np
import os
import platform
import pyttsx3
import speech_recognition as sr
import random
import string
import subprocess
import tempfile
import threading
import time
import traceback

from PIL import Image, ImageDraw
from collections import deque
from flask import Flask, render_template, redirect, url_for, request
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from screeninfo import get_monitors
from tkinter import Tk, filedialog
from typing import Callable

EXEC_PLATFORM = platform.system()
config = configparser.ConfigParser()
config.read("config.ini")
USE_GOOGLE = int(config["misc"]["use_google"])
FULLSCREEN = int(config["misc"]["fullscreen"])
DEBUG_TEST = False

# spin up ollama first, so that the reasoning portion is quicker
os.environ["OLLAMA_USE_CUDA"] = "1"
command = ["ollama", "run", "--keepalive", "10h", "mistral"]
process = subprocess.Popen(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
stdout, stderr = process.communicate()

if EXEC_PLATFORM == "Windows":
    import pygetwindow as gw

    # audio/volume control
    from ctypes import cast, POINTER
    from comtypes import CLSCTX_ALL
    from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume

app = Flask(__name__)

# Width ratios for each panel (default: equal)
width_ratios = [1, 1, 1, 1]  # Can be adjusted dynamically if needed
canvas_width = 3064
canvas_height = 672
total_ratio = sum(width_ratios)
panel_widths = [int(canvas_width * (r / total_ratio)) for r in width_ratios]
panel_widths.append(canvas_width)  # for the exp wall

# In-memory log storage
log_storage = deque(maxlen=20)
log_lock = threading.Lock()

# Configure logging
logging.basicConfig(level=logging.INFO)
voice_logger = logging.getLogger("voice_assistant")

# Sample codes for testing
random_commands = [
"Can you play the GCA video on screen one?",
"Turn on the experience wall, please.",
"I'd like to stop the video that's playing on screen four.",
"Hey, pause whatever is running on screen two.",
"Can you start the pi tech video on screen three?",
"Shut off the experience wall now.",
"Play the aero video on the second screen.",
"Lower the volume to about 30%.",
"I'd like to see the GAT video on screen 4.",
"Can you play the GAT video on screen 3?",
"Could you queue up the pi tech one on screen number two?",
"Pause the video on screen three, please.",
"Start playing the GAT video on screen 1.",
"Let’s watch cw aero on screen four.",
"Turn off the wall display thing.",
"Can you raise the volume to like 80%?",
"Stop whatever’s on screen two.",
"Play the GCA clip on the third screen.",
"I'd like to see the tech video on screen 1 — the pi one.",
        ]

def format_log_message(level, message):
    timestamp = time.strftime("%Y-%m-%d %H:%M:%S")  # Format the timestamp
    return f"[{timestamp}] [{level}] *{message}"


def log_to_memory(message, level="INFO"):
    formatted_message = format_log_message(level, message)
    with log_lock:
        try:
            previous_msg = log_storage[-1].split("*")[-1]
            if previous_msg == message:
                return
        except IndexError:
            pass  # when launched, the deque is empty
        log_storage.append(formatted_message)


def write_logs_to_file():
    log_file = "voice_assistant.log"
    while True:
        with log_lock:
            with open(log_file, "a", encoding="utf-8", errors="ignore") as f:
                for line in log_storage:
                    f.write(line + "\n")
        time.sleep(60)  # Write logs to file every 60 seconds


log_thread = threading.Thread(target=write_logs_to_file, daemon=True)
log_thread.start()


class MediaPanel:
    def __init__(self, panel_id, width):
        self.panel_id = panel_id
        self.filepath = ""
        self.cap = None
        self.playing = False
        self.video_paused = False
        self.frame = np.zeros((canvas_height, width, 3), dtype=np.uint8)
        self.width = width
        self.lock = threading.Lock()

        self.slides = []
        self.slide_index = 0
        self.slide_duration = 5
        self.fade_duration = 1
        self.last_slide_time = time.time()

    def set_video(self, path):
        with self.lock:
            self.filepath = path
            if self.cap:
                self.cap.release()
            self.cap = cv2.VideoCapture(path)
            self.playing = False
            self.fps = self.cap.get(cv2.CAP_PROP_FPS)
            self.wait_time = 1.0 / self.fps  # count time in seconds
            self.last_read_time = 0

    def play(self):
        with self.lock:
            self.playing = True

    def pause(self):
        with self.lock:
            self.playing = False
            self.video_paused = True

    def stop(self):
        with self.lock:
            self.playing = False
            self.frame = None
            if self.cap:
                self.cap.set(cv2.CAP_PROP_POS_FRAMES, 0)

    def clear_slides(self):
        with self.lock:
            self.slides = []

    def update_frame(
        self,
        start_time,
    ):
        with self.lock:
            if self.slides:
                now = time.time()
                if now - self.last_slide_time > self.slide_duration:
                    # Update slide index and last slide time
                    self.last_slide_time = now
                    next_slide_index = (self.slide_index + 1) % len(self.slides)

                    # Load current and next images
                    current_img = cv2.imread(self.slides[self.slide_index])
                    # next_img = cv2.imread(self.slides[next_slide_index])

                    # Convert images to BGRA
                    # current_rgba = cv2.cvtColor(current_img, cv2.COLOR_BGR2BGRA)
                    # next_rgba = cv2.cvtColor(next_img, cv2.COLOR_BGR2BGRA)

                    # Resize images to fit the canvas
                    current_rgba = cv2.resize(current_img, (self.width, canvas_height))
                    # next_rgba = cv2.resize(next_rgba, (self.width, canvas_height))

                    # Fade out current image and fade in next image
                    #for alpha in np.linspace(
                    #    1, 0, int(self.fade_duration * 30)
                    #):  # Assuming 30 FPS
                    #    # Create a blended image
                    #    blended = cv2.addWeighted(
                    #        current_rgba, alpha, next_rgba, 1 - alpha, 0
                    #    )
                    #    self.frame = blended
                    self.frame = current_rgba

                    # Update slide index after the transition
                    self.slide_index = next_slide_index

                else:
                    # Display the current slide
                    img = cv2.imread(self.slides[self.slide_index])
                    # rgba_image = cv2.cvtColor(img, cv2.COLOR_BGR2BGRA)
                    self.frame = cv2.resize(img, (self.width, canvas_height))

            elif self.playing and self.cap and self.cap.isOpened():
                current_time = time.time() - start_time
                if current_time - self.last_read_time >= self.wait_time:
                    ret, frame = self.cap.read()
                    if not ret:
                        # if finish playing, loop it
                        self.cap.set(cv2.CAP_PROP_POS_FRAMES, 0)
                        ret, frame = self.cap.read()
                        self.last_read_time = current_time
                    if ret:
                        if (
                            frame.shape[1] != self.width
                            or frame.shape[0] != canvas_height
                        ):
                            frame = cv2.resize(frame, (self.width, canvas_height))
                        self.frame = frame
                        self.last_read_time = current_time
                    else:
                        self.frame = None

            elif self.video_paused:
                self.frame = self.frame

            else:
                self.frame = None

    def set_slides(self, image_paths, duration=5):
        self.slides = image_paths
        self.slide_index = 0
        self.slide_duration = duration
        self.last_slide_time = time.time()


class AuthorizationCheck(dspy.Signature):
    """
    Determine if the dictated_voice_input contains or phonetically resembles an authorization code from the authorization_code_and_user dictionary.
    
    - Match is based only on the authorization code (e.g., 'E0002'), not the user name.
    - The voice input may include extra words or phrases, and may contain phonetic or transcription errors (e.g., "oh" for "0", "too" for "2", etc.).
    
    Input:
    - dictated_voice_input (str): A spoken/transcribed phrase which may loosely mention an authorization code (e.g., "Hi, I am E zero zero zero two").
    - authorization_code_and_user (dict): A mapping of known authorization codes to user names, e.g., {'E0002': 'Mr Benjamin'}.

    Output:
    - return_message (str): 
        If a match is found: "Welcome {user_name}, how can I help you?"
        If no match: "Authorization failed."
    - return_value (bool): True if a valid authorization code was found in the input; False otherwise.

    [Matching Notes]
    - Normalize common speech-to-text errors (e.g., "oh" → "0", "too" → "2", "won" → "1", etc.).
    - Ignore punctuation and extra words (e.g., "Hello, I am E zero zero three" → "E0003").
    - Use fuzzy or phonetic similarity where needed.

    Example Input:
    dictated_voice_input = "Hi, I am E zero zero two"
    authorization_code_and_user = {
        "E0001": "Mr Azib",
        "E0002": "Mr Benjamin",
        "E0003": "Mr Chan"
    }

    Expected Output:
    return_message = "Welcome Mr Benjamin, how can I help you?"
    return_value = True
    """

    dictated_voice_input: str = dspy.InputField()
    authorization_code_and_user: dict = dspy.InputField(
        desc="A dictionary with the format {authorization_code: user_name}"
    )
    return_message: str = dspy.OutputField(
        desc="A response message indicating whether access was granted or denied."
    )
    return_value: bool = dspy.OutputField(
        desc="True if the voice input matched a valid authorization code; otherwise False."
    )


class CommandControl(dspy.Signature):
    """Determine the command details based on the input request and intention.

    Input:
    - request: actual request from the user
    Output:
    # IMPORTANT:
    # Here are the valid responses:
    - play_video(video_name: str, screen_number :int)
    - pause_video(screen_number: int)
    - stop_video(screen_number: int)
    - experience_wall_control(action: str)  # valid actions: on/off
    - set_volume(volume_level: float)  # 0.0 to 1.0 value for volume_level
    - system_sleep()

    [PRIORITY] Try your best to match file names and screen numbers.
    [FALLBACK] If there's no matching file names, try using similar sounding phonetics to match the valid video files names below.

    # Video Controls:
    ## Valid video files names:
    gca_video
    cw_aero_video
    pi_tech_video
    gat_video

    # Valid screen number are (1,2,3,4). However, please map it to (0, 1, 2, 3)

    ## Output format:
    The output should only contain the function name and positional args.
    """

    request: str = dspy.InputField()

    command: str = dspy.OutputField(desc="Actual function and parameters to run")
    confidence: float = dspy.OutputField(
        desc="For downstream error handling or routing"
    )  # For downstream error handling or routing


# Initialize 4 panels + exp wall panel
media_panels = [MediaPanel(panel_id=i, width=panel_widths[i]) for i in range(5)]
experience_wall = media_panels[-1]


# Background thread: single OpenCV window to display all panels
def display_loop():
    cv2.namedWindow("Media Dashboard", cv2.WINDOW_NORMAL)
    cv2.resizeWindow("Media Dashboard", canvas_width, canvas_height)

    # Move display window to the target monitor
    # Function to get the position of the monitor with the specified name
    def get_monitor_position(monitor_name):
        for monitor in get_monitors():
            print(f"DEBUG: {monitor=}, {monitor_name=}, {monitor.name == monitor_name}")
            if monitor_name == monitor.name:
                return (monitor.x, monitor.y, monitor.width, monitor.height)
        return None

    # Get the position of the HDMI-1 monitor
    config = configparser.ConfigParser()
    config.read("config.ini")

    print(config["screen"]["target"])
    target_monitor = get_monitor_position(config["screen"]["target"])

    if target_monitor is None:
        print("target monitor not found.")
    else:
        # Get the window with the title "abc"
        try:
            window = gw.getWindowsWithTitle("Media Dashboard")[
                0
            ]  # Get the first window with the title

            # Move the window to the HDMI-1 monitor
            window.moveTo(target_monitor[0], target_monitor[1])
            window.maximize()

            print(f"Moved window '{window.title}' to HDMI-1 monitor.")
        except IndexError:
            print("Window with title 'Media Dashboard' not found.")

    # Function to overlay frame2 over frame1
    def overlay_frames(frame1, frame2):
        # Ensure both frames are the same size
        if frame1.shape != frame2.shape:
            raise ValueError("Frames must be the same size")

        # Create a copy of frame1 to modify
        combined_frame = frame1.copy()

        # Check where frame1 is transparent
        alpha1 = frame1[:, :, 3]  # Get the alpha channel of frame1
        mask = alpha1 == 0  # Create a mask where frame1 is transparent

        # Overlay frame2 where frame1 is transparent
        combined_frame[mask] = frame2[
            mask
        ]  # Set pixels from frame2 where frame1 is transparent

        return combined_frame

    start_time = time.time()
    while True:
        # first, we fill it with the experience wall content
        experience_wall.update_frame(start_time)
        exp_frames = experience_wall.frame
        if exp_frames is None:
            exp_frames = np.zeros(
                (canvas_height, experience_wall.width, 3), dtype=np.uint8
            )
        backup_frame = copy.deepcopy(exp_frames)

        # just to make sure the frames are same size
        if exp_frames.shape[1] != canvas_width or exp_frames.shape[0] != canvas_height:
            exp_frames = cv2.resize(exp_frames, (canvas_width, canvas_height))

        for i, panel in enumerate(media_panels[:-1]):
            x_start = i * panel.width
            x_end = x_start + panel.width

            panel.update_frame(start_time)
            if panel.frame is not None:
                exp_frames[:, x_start:x_end] = panel.frame

        cv2.imshow("Media Dashboard", exp_frames)
        if FULLSCREEN:
            cv2.setWindowProperty(
                "Media Dashboard", cv2.WND_PROP_FULLSCREEN, cv2.WINDOW_FULLSCREEN
            )
        # Check for 'Ctrl + Q' key press
        if cv2.waitKey(1) & 0xFF == ord("q"):
            break
    cv2.destroyAllWindows()

threading.Thread(target=display_loop, daemon=True).start()


def log_function_call(func):
    def wrapper(*args, **kwargs):
        # Log the function name and arguments
        print(f"Calling function: {func.__name__}")
        print(f"Arguments: {args}, Keyword Arguments: {kwargs}")

        # Call the original function
        result = func(*args, **kwargs)

        # Optionally log the result
        print(f"Function {func.__name__} returned: {result}")

        return result

    return wrapper


def remove_punctuation(text):
    translator = str.maketrans("", "", string.punctuation)
    return text.translate(translator)


@log_function_call
def play_video(video_name="", panel_index=0, *args, **kwargs):
    config = configparser.ConfigParser()
    config.read("config.ini")

    video_dict = config["presentation"]
    if video_name:
        media_panels[panel_index].set_video(video_dict[video_name])

    if panel_index == "all":
        for p in media_panels:
            p.set_video(video_dict[video_name])
            p.play()
    else:
        media_panels[panel_index].play()


def pause_video(panel_index=0, *args, **kwargs):
    if panel_index == "all":
        for p in media_panels:
            p.pause()
    else:
        media_panels[panel_index].pause()


def stop_video(panel_index=0, *args, **kwargs):
    if panel_index == "all":
        for p in media_panels:
            p.stop()
    else:
        media_panels[panel_index].stop()
        exp_wall_paused = media_panels[-1].video_paused
        if exp_wall_paused:
            media_panels[-1].play()
            time.sleep(0.2)
            media_panels[-1].pause()


def experience_wall_control(action, *args, **kwargs):
    if action.lower() == "on":
        play_video(video_name="experience_wall", panel_index=-1)
    elif action.lower() == "off":
        stop_video(panel_index=-1)


def set_volume(volume_level, *args, **kwargs):
    # Ensure volume_level is between 0.0 and 1.0
    if not (0.0 <= volume_level <= 1.0):
        raise ValueError("Volume level must be between 0.0 and 1.0")

    # Get the audio devices
    speaker_devices = AudioUtilities.GetSpeakers()

    # Get the volume interface for the default audio device
    speaker_interface = speaker_devices.Activate(
        IAudioEndpointVolume._iid_, CLSCTX_ALL, None
    )

    # Cast to IAudioEndpointVolume to manipulate the volume
    volume = cast(speaker_interface, POINTER(IAudioEndpointVolume))

    # Set the master volume level
    volume.SetMasterVolumeLevelScalar(volume_level, None)


def system_sleep(*args, **kwargs):
    for p in media_panels:
        p.stop()


class VoiceAssistant:
    def __init__(self, log_to_memory, command_registry):
        self.log_to_memory = log_to_memory
        self.config = configparser.ConfigParser()
        self.config.read("config.ini")

        # DSPY related
        ollama_server_url = self.config["llm"]["endpoint_prod"]
        self.lm = dspy.LM(
            self.config["llm"]["model"], api_base=ollama_server_url, api_key=""
        )
        dspy.configure(lm=self.lm)
        self.get_command_actions = dspy.ChainOfThought(CommandControl)
        self.authorization_check = dspy.ChainOfThought(AuthorizationCheck)
        self.verified = False

        # voice engine
        self.voice_engine = pyttsx3.init("sapi5")
        self.voices = self.voice_engine.getProperty("voices")
        self.voice_engine.setProperty(
            "voice", self.voices[1].id
        )  # changing index, changes voices. 1 for female

        # dictionary of video files
        self.video_dict = self.config["presentation"]
        self.command_registry = command_registry

        self.recognizer = sr.Recognizer()
        self.mic_index = int(self.config["voice_control"]["mic_index"])

    def speak(self, audio):
        self.voice_engine.say(audio)
        self.voice_engine.runAndWait()

    def execute_command(self, command_str: str):
        try:
            # Safely parse the command string
            tree = ast.parse(command_str, mode="eval")
            if not isinstance(tree.body, ast.Call):
                raise ValueError("Command must be a function call")

            func_name = tree.body.func.id

            args = []
            for a in tree.body.args:
                args.append(a.value)

            kargs = {}
            for kw in tree.body.keywords:
                kargs[kw.arg] = ast.literal_eval(kw.value)

            if func_name not in self.command_registry:
                raise ValueError(f"Unknown command: {func_name}")

            self.command_registry[func_name](*args, **kargs)

        except Exception as e:
            print(f"Failed to execute command: {e}, Received: {command_str=}")
            self.log_to_memory(traceback.print_exc(), level="DEBUG")
            print(traceback.print_exc())

    def verify_user(self, source):
        if not int(self.config["misc"]["use_authorization"]):
            return True

        self.speak("Hello, please provide identification code.")
        self.log_to_memory("Hello, please provide identification code.")
        audio = self.recognizer.listen(source)
        if USE_GOOGLE:
            transcript = self.recognizer.recognize_google(audio)
        else:
            transcript = self.recognizer.recognize_whisper(audio, model='medium.en')

        if DEBUG_TEST:
            transcript = 'E0001'

        self.log_to_memory(f'you said: {transcript}')
        res = self.authorization_check(
            dictated_voice_input=transcript,
            authorization_code_and_user=self.config["authorization"],
        )
        self.log_to_memory(f"Authorization check results: {res}")
        message = res.return_message
        value = res.return_value
        self.speak(message)
        if value:
            self.verified = True
        return value

    def unauthorized_access(self):
        self.speak("False identification, I do not recognize you.")

    def listen_for_wake_phrase(self, source):
        self.log_to_memory("Listening for wake phrases")
        wake_phrases = [
            "hello aeries",
            "hello aries",
            "hello iris",
            "hello aeris",
            "hello i miss",
            "hello irene",
            "hello ivy",
            "hello i mean",
            "hello im reese",
            "hello ill reach",
            "hello elise",
        ]
        while True:
            try:
                audio = self.recognizer.listen(source)
                if USE_GOOGLE:
                    transcript = self.recognizer.recognize_google(audio)
                else:
                    transcript = self.recognizer.recognize_whisper(audio, model='medium.en')

                transcript = remove_punctuation(transcript)

                if DEBUG_TEST:
                    transcript = 'hello aeris'
                if self.is_valid_command(transcript):
                    self.log_to_memory(f"You said: {transcript}", level="INFO")

                if any(w.lower() in transcript.lower() for w in wake_phrases):
                    self.log_to_memory(
                        "Wake word detected. Listening for command...", level="INFO"
                    )
                    self.try_to_recognize(source)
            except sr.UnknownValueError:
                self.log_to_memory("Didn't catch any sound.", level="INFO")
            except sr.RequestError as e:
                self.log_to_memory(f"Recognition error: {e}", level="ERROR")
                self.speak("Recognition error.")
                print(traceback.print_exc())
            except Exception as e:
                self.log_to_memory(f"Unhandled error: {e}", level="ERROR")
                self.log_to_memory(traceback.print_exc(), level="DEBUG")

    def is_valid_command(self, message):
        """Check if the log message is valid."""
        return message and message.strip() not in ["", "...", "... ... ... ..."]

    def try_to_recognize(self, source, repeated=False):
        if not self.verified:
            if not self.verify_user(source):
                return

        if not repeated:
            self.speak("Hello, you have called for me. How can i help you?")

        self.log_to_memory("You may speak your command now.")
        audio = self.recognizer.listen(
            source,
        )
        if USE_GOOGLE:
            command = self.recognizer.recognize_google(audio)
        else:
            command = self.recognizer.recognize_whisper(audio, model='medium.en')

        if DEBUG_TEST:
            command = random.choice(random_commands)

        self.log_to_memory(f"Command received: {command}")

        if self.is_valid_command(command):
            command_res = self.get_command_actions(
                request=command,
            )
            if command_res.confidence <= float(
                self.config["voice_control"]["confidence"]
            ):
                self.speak("I'm sorry, could you repeat please?")
                self.try_to_recognize(source, repeated=True)

            self.speak(f"Okay running: {command}")
            self.log_to_memory(f"Generated command: {command_res}")
            # if it passes the conf threshold
            self.execute_command(command_res.command)


    def run(self):
        try:
            with sr.Microphone(device_index=self.mic_index) as source:
                self.log_to_memory("Calibrating for ambient noise...")
                self.recognizer.adjust_for_ambient_noise(source, duration=2)
                self.recognizer.pause_threshold = (
                    2  # add some extra time to allow brain to process more
                )
                self.log_to_memory("Always listening for wake word...")
                self.listen_for_wake_phrase(source)
        except Exception as mic_error:
            self.log_to_memory(
                f"Microphone initialization failed: {mic_error}", level="ERROR"
            )
            self.log_to_memory(traceback.print_exc(), level="DEBUG")


def run_aeris():
    log_to_memory("")
    log_to_memory(" ---------------- ")
    log_to_memory(" STARTING AERIS SERVICE")
    log_to_memory(" ---------------- ")
    log_to_memory("")

    for index, name in enumerate(sr.Microphone.list_microphone_names()):
        print(f"Microphone {index}: {name}")

    # if mic_index is None:
    #     mic_index = input("Which mic would you like to use? :") or 8
    #     mic_index = int(mic_index)
    #
    command_registry: dict[str, Callable] = {
        "play_video": play_video,
        "pause_video": pause_video,
        "stop_video": stop_video,
        "experience_wall_control": experience_wall_control,
        "set_volume": set_volume,
        "system_sleep": system_sleep,
    }

    aeris = VoiceAssistant(
        log_to_memory=log_to_memory, command_registry=command_registry
    )
    aeris.run()


# if there's powerpoint slide
def convert_pptx_to_images(pptx_path: str) -> list[str]:
    output_dir = Path(tempfile.mkdtemp())
    ppt = comtypes.client.CreateObject("PowerPoint.Application")
    ppt.Visible = 1
    presentation = ppt.Presentations.Open(pptx_path, WithWindow=False)

    export_path = str(output_dir / "slide")
    presentation.SaveAs(export_path, 17)  # 17 = ppSaveAsJPG
    presentation.Close()
    ppt.Quit()

    # Collect all slide image paths
    slide_images = sorted(output_dir.glob("*.JPG"))
    return [str(p) for p in slide_images]


def convert_pptx_to_images_pure(pptx_path, width=1280, height=720) -> list[str]:
    prs = Presentation(pptx_path)
    output_dir = Path(tempfile.mkdtemp())

    slide_images = []
    for idx, slide in enumerate(prs.slides):
        img = Image.new("RGB", (width, height), "white")
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.has_text_frame:
                if shape.text_frame:
                    text = shape.text_frame.text
                    left = int(shape.left * width / prs.slide_width)
                    top = int(shape.top * height / prs.slide_height)
                    draw.text((left, top), text, fill="black")  # Add font if needed

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_stream = shape.image.blob
                with Image.open(io.BytesIO(img_stream)) as shape_img:
                    shape_img = shape_img.convert("RGB")
                    shape_img = shape_img.resize(
                        (shape.width // 9525, shape.height // 9525)
                    )
                    img.paste(shape_img, (shape.left // 9525, shape.top // 9525))

        image_path = output_dir / f"slide_{idx + 1}.jpg"
        img.save(image_path)
        slide_images.append(str(image_path))

    return slide_images


@app.route("/")
def index():
    config = configparser.ConfigParser()
    return render_template("dashboard.html", panels=media_panels, config=config)


@app.route("/panel/<int:panel_id>/browse", methods=["POST"])
def browse_file(panel_id):
    print(f"Browse clicked for panel {panel_id}")
    root = Tk()
    root.withdraw()
    root.attributes("-topmost", True)
    file_path = filedialog.askopenfilename(
        filetypes=[("Video files", "*.mp4 *.avi *.mkv *.mov"), ("All files", "*.*")]
    )
    root.destroy()
    print(f"Selected file: {file_path}")
    if file_path:
        media_panels[panel_id].set_video(file_path)
    return redirect(url_for("index"))


@app.route("/panel/<int:panel_id>/<action>", methods=["POST"])
def control_panel(panel_id, action):
    panel = media_panels[panel_id]
    if action == "play":
        panel.play()
    elif action == "pause":
        panel.pause()
    elif action == "stop":
        panel.stop()
        exp_wall_paused = media_panels[-1].video_paused
        if exp_wall_paused:
            media_panels[-1].play()
            time.sleep(0.2)
            media_panels[-1].pause()
    elif action == "clear_slides":
        panel.clear_slides()

    return redirect(url_for("index"))


@app.route("/bulk_action", methods=["POST"])
def bulk_action():
    action = request.form.get("action")
    panel_ids = request.form.getlist("panel_ids")  # this is a list of strings

    for panel_id_str in panel_ids:
        try:
            panel_id = int(panel_id_str)
            panel = media_panels[panel_id]
            if action == "play":
                panel.play()
            elif action == "pause":
                panel.pause()
            elif action == "stop":
                panel.stop()
            elif action == "clear_slides":
                panel.clear_slides()
        except (ValueError, IndexError):
            # Handle invalid panel ID
            continue

    return redirect(url_for("index"))


@app.route("/panel/<int:panel_id>/upload_pptx", methods=["POST"])
def upload_pptx(panel_id):
    root = Tk()
    root.withdraw()
    root.attributes("-topmost", True)

    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    root.destroy()

    if file_path:
        # slide_images = convert_pptx_to_images_pure(file_path)
        slide_images = convert_pptx_to_images(file_path)
        media_panels[panel_id].set_slides(
            slide_images, duration=5
        )  # set default time here

    return redirect(url_for("index"))


@app.route("/panel/<int:panel_id>/upload_images", methods=["POST"])
def upload_images(panel_id):
    root = Tk()
    root.withdraw()
    root.attributes("-topmost", True)

    folder_path = filedialog.askdirectory()  # Ask for a directory instead of a file
    root.destroy()

    if folder_path:
        # Get a list of image files in the selected directory
        image_extensions = (
            ".png",
            ".jpg",
            ".jpeg",
            ".gif",
            ".bmp",
            ".webp",
        )  # Add more extensions if needed
        slide_images = [
            os.path.join(folder_path, f)
            for f in os.listdir(folder_path)
            if f.lower().endswith(image_extensions)
        ]

        media_panels[panel_id].set_slides(
            slide_images, duration=5
        )  # Set default time here

    return redirect(url_for("index"))


@app.route("/logs")
def view_logs():
    try:
        with log_lock:
            # Get the last num_lines from memory
            lines = log_storage
            # Ensure all lines are strings
            lines = [line for line in lines if isinstance(line, str)]
    except Exception as e:
        lines = [f"Error reading log: {e}"]

    # Join the lines safely
    return "<br>".join(line.strip() for line in lines)


if __name__ == "__main__":
    assistant_thread = threading.Thread(
        target=run_aeris,
        daemon=True,
    )
    assistant_thread.start()

    app.run(debug=False, use_reloader=False)
